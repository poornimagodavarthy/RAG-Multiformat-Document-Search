import shutil
from pathlib import Path
import json
from datetime import datetime
import re

import pymupdf4llm
import pandas as pd
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from api.database import SessionLocal, DocumentMetadata
from ingestion.chunking import check_document_in_vectordb

import boto3
from botocore.exceptions import ClientError
import os
import uuid
from api.config import ClientConfig
import tempfile

# Initialize S3 client
s3_client = boto3.client(
    's3',
    region_name=ClientConfig.S3_REGION,
    aws_access_key_id=ClientConfig.AWS_ACCESS_KEY_ID,
    aws_secret_access_key=ClientConfig.AWS_SECRET_ACCESS_KEY
)

# S3 paths
S3_PREFIX_ORIGINAL = "original_documents/"
S3_PREFIX_PARSED = "parsed_markdowns/"

# PATH SETUP
BASE_DIR = Path.cwd()
PARSED_DIR = BASE_DIR / "parsed_data"
PARSED_DIR.mkdir(exist_ok=True)

def upload_to_s3(file_path: Path, s3_key: str) -> str:
    """
    Upload a file to S3 and return the URL.
    """
    try:
        with open(file_path, 'rb') as f:
            s3_client.upload_fileobj(f, ClientConfig.S3_BUCKET, s3_key)  
        s3_url = f"https://{ClientConfig.S3_BUCKET}.s3.{ClientConfig.S3_REGION}.amazonaws.com/{s3_key}"
        print(f"[S3] Uploaded {file_path.name} → {s3_url}")
        return s3_url
    except ClientError as e:
        print(f"[S3 ERROR] Failed to upload {file_path}: {e}")
        return None
    except Exception as e:
        print(f"[S3 ERROR] Failed to open or upload {file_path}: {e}")
        return None
    
def upload_original_document(file_path: Path) -> str:
    """Upload original document to S3 original_documents/ folder"""
    s3_key = f"{S3_PREFIX_ORIGINAL}{file_path.name}"
    return upload_to_s3(file_path, s3_key)

def upload_parsed_markdown(file_path: Path) -> str:
    """Upload parsed markdown to S3 parsed_markdowns/ folder"""
    s3_key = f"{S3_PREFIX_PARSED}{file_path.name}"
    return upload_to_s3(file_path, s3_key)

def convert_markdown_to_pdf(
    md_content: str,
    output_path: Path,
    title: str = "",
    original_filename: str = ""
) -> bool:
    """
    Convert markdown to a neutral PDF for rendering.
    """
    try:
        import re
        import markdown
        from weasyprint import HTML

        # Remove ingestion page markers
        md_clean = re.sub(r'<!-- PAGE \d+ -->', '', md_content)

        # Markdown → HTML
        html_body = markdown.markdown(
            md_clean,
            extensions=[
                "extra",
                "tables",
                "fenced_code",
                "sane_lists",
            ]
        )

        html = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="utf-8" />
            <title>{title}</title>

            <style>
                @page {{
                    size: letter;
                    margin: 1in allowing;

                    @bottom-center {{
                        content: "Page " counter(page) " of " counter(pages);
                        font-size: 9pt;
                        color: #000;
                    }}
                }}

                body {{
                    font-family: Arial, Helvetica, sans-serif;
                    font-size: 12.5px;
                    line-height: 1.65;
                    color: #000;
                }}

                p {{
                    margin: 0 0 0.9em 0;
                    color: #000;
                }}

                /* Headings — pure Google Docs hierarchy */
                h1 {{
                    font-size: 22px;
                    font-weight: 600;
                    margin: 24px 0 10px;
                    color: #000;
                }}

                h2 {{
                    font-size: 17px;
                    font-weight: 600;
                    margin: 20px 0 8px;
                    color: #000;
                }}

                h3 {{
                    font-size: 15px;
                    font-weight: 600;
                    margin: 16px 0 6px;
                    color: #000;
                }}

                h4 {{
                    font-size: 14px;
                    font-weight: 600;
                    margin: 14px 0 6px;
                    color: #000;
                }}

                ul, ol {{
                    padding-left: 18px;
                    margin: 8px 0;
                }}

                li {{
                    margin: 4px 0;
                    color: #000;
                }}

                table {{
                    width: 100%;
                    border-collapse: collapse;
                    margin: 12px 0;
                    font-size: 12px;
                    color: #000;
                }}

                th {{
                    text-align: left;
                    font-weight: 600;
                    padding: 6px 8px;
                    border-bottom: 1px solid #ccc;
                    color: #000;
                }}

                td {{
                    padding: 6px 8px;
                    border-bottom: 1px solid #e0e0e0;
                    vertical-align: top;
                    color: #000;
                }}

                blockquote {{
                    margin: 12px 0;
                    padding-left: 12px;
                    border-left: 3px solid #ccc;
                    color: #000;
                }}

                code {{
                    background: #f5f5f5;
                    padding: 2px 4px;
                    font-size: 11px;
                    color: #000;
                }}

                pre {{
                    background: #f5f5f5;
                    padding: 10px;
                    margin: 12px 0;
                    font-size: 11px;
                    white-space: pre-wrap;
                    color: #000;
                }}

                a {{
                    color: #000;
                    text-decoration: underline;
                }}

                h1, h2, h3 {{
                    page-break-after: avoid;
                }}

                li {{
                    page-break-inside: avoid;
                }}
            </style>
        </head>

        <body>
            {html_body}
        </body>
        </html>
        """

        HTML(string=html).write_pdf(str(output_path))
        return True

    except Exception as e:
        print(f"[PDF CONVERT ERROR] {e}")
        return False

    
# METADATA HELPERS
def save_metadata_json_backup():
    """Export database metadata to JSON for debugging/backup"""
    db = SessionLocal()
    try:
        docs = db.query(DocumentMetadata).all()
        metadata = {}
        for doc in docs:
            metadata[doc.markdown_filename] = {
                "document_id": doc.document_id,
                "title": doc.title,
                "original_filename": doc.original_filename,
                "markdown_filename": doc.markdown_filename,
                "type": doc.type,
                "original_s3_url": doc.original_s3_url, 
                "parsed_s3_url": doc.parsed_s3_url,    
                "date_updated": doc.date_updated,
                "total_pages": doc.total_pages,
                "file_size": doc.file_size
            }
        
        # Save to JSON file
        backup_path = PARSED_DIR / "metadata_backup.json"
        backup_path.write_text(json.dumps(metadata, indent=2))
        print(f"[BACKUP] Metadata saved to {backup_path}")
    finally:
        db.close()

def save_document_metadata(markdown_filename, metadata_dict, client_id: str):
    """Save or update document metadata in database"""
    db = SessionLocal()
    try:
        # Check if exists
        doc = db.query(DocumentMetadata).filter(
            DocumentMetadata.markdown_filename == markdown_filename,
            DocumentMetadata.client_id == client_id
        ).first()
        
        if doc:
            for key, value in metadata_dict.items():
                setattr(doc, key, value)
        else:
            # Create new
            doc = DocumentMetadata(
                client_id=client_id,
                markdown_filename=markdown_filename,
                **metadata_dict
            )
            db.add(doc)
        
        db.commit()
    finally:
        db.close()


def create_and_save_metadata(file_path: Path, dst: Path, total_pages: int, client_id: str, content: str = None):
    """
    Helper to create and save metadata for any processed document.
    Converts text-based formats to PDF for viewing.
    """
    stat = file_path.stat()
    file_size_mb = round(stat.st_size / 1_000_000, 2)
    updated_date = datetime.fromtimestamp(stat.st_mtime).strftime("%b %Y")
    db = SessionLocal()

    # Detect doc type
    doc_type = detect_doc_type(file_path, content)
    
    # Create clean title from filename
    title = create_clean_title(file_path.stem)
    
    # Determine if we need to convert to PDF for viewing
    file_ext = file_path.suffix.lower()
    needs_pdf_conversion = file_ext in ['.ppt', '.pptx', '.txt', '.md', '.markdown']
    is_spreadsheet = file_ext in ['.xlsx', '.xls', '.csv']
    
    # Upload original file first (for downloading)
    download_url = upload_original_document(file_path)
    
    # Initialize to avoid UnboundLocalError
    parsed_s3_url = None

    # Upload original file first (for downloading)
    download_url = upload_original_document(file_path)

    if needs_pdf_conversion and content:
        # Convert text/markdown → PDF for viewing
        import tempfile
        
        # Create PDF in temp directory
        with tempfile.NamedTemporaryFile(mode='wb', suffix='_view.pdf', delete=False) as tmp_file:
            pdf_path = Path(tmp_file.name)
        
        try:
            is_markdown = file_ext in ['.md', '.markdown']

            if convert_markdown_to_pdf(content, pdf_path, title):
                pdf_s3_key = f"{S3_PREFIX_ORIGINAL}{file_path.stem}_view.pdf"
                view_url = upload_to_s3(pdf_path, pdf_s3_key)
            else:
                view_url = download_url
        finally:
            # Clean up temp PDF file
            pdf_path.unlink(missing_ok=True)

        # Parsed markdown is always uploaded
        parsed_s3_url = upload_parsed_markdown(dst)

    elif is_spreadsheet:
        # Excel / CSV to CSV
        parsed_s3_url = upload_parsed_markdown(dst)
        view_url = parsed_s3_url

    else:
        # PDFs
        view_url = download_url
        parsed_s3_url = upload_parsed_markdown(dst)
    
    # Generate unique document ID
    existing = db.query(DocumentMetadata).filter(
        DocumentMetadata.client_id == client_id,
        DocumentMetadata.markdown_filename == dst.name
    ).first()

    if existing:
        document_id = existing.document_id
    else:
        document_id = str(uuid.uuid4())

    
    metadata_entry = {
        "document_id": document_id,
        "title": title,
        "original_filename": file_path.name,
        "type": doc_type,
        "original_s3_url": view_url, # For viewing (CSV for Excel, PDF for DOCX, etc.)
        "download_s3_url": download_url,  # For downloading (original file)
        "parsed_s3_url": parsed_s3_url,
        "date_updated": updated_date,
        "total_pages": total_pages,
        "file_size": file_size_mb
    }
    
    save_document_metadata(dst.name, metadata_entry, client_id)
    
    print(f"[METADATA SAVED] for {dst.name} (Type: {doc_type}, ID: {document_id})")
    print(f"  View URL: {view_url}")
    print(f"  Download URL: {download_url}")
def create_clean_title(stem: str) -> str:
    """
    Convert a filename stem to a clean, human-readable title.
    """
    # Replace underscores and hyphens with spaces
    title = stem.replace('_', ' ').replace('-', ' ')
    
    # Remove leading numbers (like '03' or '07')
    title = re.sub(r'^\d+\s+', '', title)
    
    # Remove common prefixes that are noise
    title = re.sub(r'^(export|notion export)\s+', '', title, flags=re.IGNORECASE)
    
    # Collapse multiple spaces
    title = re.sub(r'\s+', ' ', title).strip()
    
    # Title case, but keep acronyms uppercase
    words = title.split()
    cleaned_words = []
    
    for word in words:
        # Keep acronyms (all caps, 2-5 letters) as-is
        if word.isupper() and 2 <= len(word) <= 5:
            cleaned_words.append(word)
        # Keep words that are all digits
        elif word.isdigit():
            cleaned_words.append(word)
        # Title case everything else
        else:
            cleaned_words.append(word.capitalize())
    
    return ' '.join(cleaned_words)

# FILE HANDLERS

def detect_doc_type(file_path: Path, content: str = None) -> str:
    """
    Detect document type using filename and optional content heuristics.
    """
    filename = file_path.name.lower()
    
    # PRIORITY CHECK: SOP anywhere in filename
    if 'sop' in filename:
        return "SOP"
    
    # Define patterns for each doc type
    patterns = {
        "Research Paper": [
            r'paper',
            r'research',
            r'study',
            r'abstract',
            r'journal',
            r'thesis',
        ],
        "Report": [
            r'report',
            r'analysis',
            r'summary',
            r'findings',
            r'review',
        ],
        "Technical Doc": [
            r'spec',
            r'specification',
            r'technical',
            r'architecture',
            r'documentation',
            r'api',
            r'manual',
            r'handbook',
        ],
        "Policy": [
            r'policy',
            r'policies',
            r'guidelines',
            r'rules',
            r'standards',
            r'compliance',
        ],
        "Presentation": [
            r'\.pptx?$',
            r'presentation',
            r'slides',
            r'deck',
        ],
        "Spreadsheet": [
            r'\.xlsx?$',
            r'\.csv$',
            r'data',
            r'sheet',
        ],
        "Notes": [
            r'notes',
            r'meeting',
            r'minutes',
            r'transcript',
            r'\.md$',
        ],
        "Contract": [
            r'contract',
            r'agreement',
            r'terms',
            r'legal',
        ],
    }
    
    # Check filename against all patterns
    for doc_type, pattern_list in patterns.items():
        for pattern in pattern_list:
            if re.search(pattern, filename, re.IGNORECASE):
                return doc_type
    
    # If content is provided, check first 500 chars for keywords
    if content:
        content_sample = content[:500].lower()
        
        content_keywords = {
            "Research Paper": ["abstract", "methodology", "conclusion", "hypothesis"],
            "Technical Doc": ["specification", "implementation", "architecture", "endpoint"],
            "Report": ["findings", "analysis", "summary", "results"],
            "Policy": ["policy", "compliance", "requirement", "standard"],
        }
        
        for doc_type, keywords in content_keywords.items():
            if sum(1 for kw in keywords if kw in content_sample) >= 2:
                return doc_type
    
    return "Unknown"

def ensure_metadata_exists(file_path: Path, dst: Path, total_pages: int, client_id: str, content: str = None):
    """Ensure a DB row exists even if markdown already exists."""
    db = SessionLocal()
    try:
        existing = db.query(DocumentMetadata).filter(
            DocumentMetadata.markdown_filename == dst.name,
            DocumentMetadata.client_id == client_id
        ).first()
    finally:
        db.close()

    if existing:
        print(f"[METADATA] Already exists for {dst.name}")
        return
    ext = file_path.suffix.lower()

    # SAFE REDIRECT FOR POWERPOINT
    if ext in ['.ppt', '.pptx']:
        print("[METADATA] PPTX missing metadata — reprocessing via handle_powerpoint")
        handle_powerpoint(file_path, client_id)
        return

    print(f"[METADATA] Missing for {dst.name}, creating now")
    create_and_save_metadata(file_path, dst, total_pages, client_id, content)

def copy_markdown(src_path: Path, dst_path: Path):
    print(f"[COPY] MD → MD: {src_path} -> {dst_path}")
    shutil.copy(src_path, dst_path)


def handle_pdf(file_path, client_id: str):
    print(f"[PDF] Client {client_id}: Processing {file_path}")

    client_parsed_dir = PARSED_DIR / client_id
    client_parsed_dir.mkdir(parents=True, exist_ok=True)
    dst = client_parsed_dir / f"{file_path.stem}.md"
    if dst.exists():
        print(f"[SKIP] Markdown already exists → {dst}")
        ensure_metadata_exists(file_path, dst, total_pages=len(fitz.open(file_path)), client_id=client_id)
        return

    pdf = fitz.open(file_path)
    md_lines = []
    
    for page_num in range(len(pdf)):
        md_lines.append(f"\n<!-- PAGE {page_num + 1} -->\n")

        try:
            page_md = pymupdf4llm.to_markdown(
                str(file_path),
                pages=[page_num],
            )

            if page_md:
                md_lines.append(page_md)
            else:
                md_lines.append("_[No extractable text on this page]_")

        except Exception as e:
            print(
                f"[PDF WARNING] Failed to parse page {page_num + 1} "
                f"of {file_path.name}: {e}"
            )
            md_lines.append("_[Page could not be parsed]_")

    
    md_content = "\n".join(md_lines)
    dst.write_text(md_content, encoding="utf-8")
    print(f"[SAVED] {dst}")

    create_and_save_metadata(file_path, dst, len(pdf), client_id, md_content)


def docx_md_helper(run):
    text = run.text
    if run.bold:
        text = f"**{text}**"
    if run.italic:
        text = f"*{text}*"
    return text
def process_single_file_from_s3(s3_key: str, filename: str, client_id: str):
    """
    Process a file from S3 through the full pipeline:
    1. Download from S3
    2. Parse document → markdown/CSV
    3. Upload parsed to S3
    4. Create SQL metadata
    5. Chunk and embed into vector DB
    """
    
    try:
        print(f"[PROCESS] Client {client_id}: Starting pipeline for {filename} from S3")
        
        # Import dependencies
        from ingestion.chunking import chunk_document, load_metadata_from_db
        
        # Step 1: Download file from S3 to temp location WITH ORIGINAL FILENAME
        temp_dir = Path(tempfile.mkdtemp())
        safe_filename = Path(filename).name
        temp_path = temp_dir / safe_filename

        try:
            print(f"[PROCESS] Downloading from S3: {s3_key}")
            s3_client.download_file(ClientConfig.S3_BUCKET, s3_key, str(temp_path))
            
            # Step 2: Process the file (this will parse and create metadata)
            handlers = get_file_handlers()
            extension = temp_path.suffix.lower()
            handler = handlers.get(extension, handle_unknown)


            if handler == handle_unknown:
                raise ValueError(f"Unsupported file type: {extension}")
            
            # Call the appropriate handler
            handler(temp_path, client_id)
            
            # Step 3: Determine the parsed filename
            original_stem = Path(filename).stem
            original_ext = Path(filename).suffix.lower()
            
            if original_ext in ['.pdf', '.docx', '.doc', '.txt', '.md', '.markdown', '.pptx', '.ppt']:
                parsed_ext = '.md'
            elif original_ext in ['.csv', '.xlsx', '.xls']:
                parsed_ext = '.csv'
            else:
                raise ValueError(f"Unsupported file type: {original_ext}")
            
            parsed_filename = f"{original_stem}{parsed_ext}"
            
            # Step 4: Load metadata from SQL
            metadata = load_metadata_from_db(client_id=client_id)
            metadata_entry = metadata.get(parsed_filename)
            
            if not metadata_entry:
                raise ValueError(f"No metadata found for {parsed_filename}")
            
            document_id = metadata_entry.get("document_id")
            if not document_id:
                raise ValueError(f"No document_id in metadata for {parsed_filename}")
            
            # Step 5: Get parsed file from S3 (it was uploaded by the handler)
            parsed_s3_url = metadata_entry.get("parsed_s3_url")
            if not parsed_s3_url:
                raise ValueError(f"No parsed_s3_url for {parsed_filename}")
            
            # Download parsed file to temp location for chunking
            parsed_s3_key = parsed_s3_url.split('.amazonaws.com/')[1]
            with tempfile.NamedTemporaryFile(delete=False, suffix=parsed_ext) as parsed_tmp:
                parsed_temp_path = Path(parsed_tmp.name)
            
            try:
                print(f"[PROCESS] Downloading parsed file from S3: {parsed_s3_key}")
                s3_client.download_file(ClientConfig.S3_BUCKET, parsed_s3_key, str(parsed_temp_path))
                
                document_id = metadata_entry.get("document_id")
                if check_document_in_vectordb(document_id, client_id):
                    print(f"[PROCESS] Document {parsed_filename} already in vector DB, skipping chunking")
                else:
                    print(f"[PROCESS] Chunking {parsed_filename}")
                    chunk_document(parsed_temp_path, metadata_entry, client_id=client_id, batch_size=32)
                
            finally:
                # Clean up parsed temp file
                parsed_temp_path.unlink(missing_ok=True)
            
        finally:
            # Clean up temp directory and all files
            try:
                if temp_dir.exists():
                    print(f"[PROCESS] Cleaning up temp directory: {temp_dir}")
                    shutil.rmtree(temp_dir)
                    print(f"[PROCESS] Cleanup complete")
            except Exception as cleanup_error:
                print(f"[PROCESS WARNING] Failed to cleanup temp directory: {cleanup_error}")
                # Don't raise - cleanup failure shouldn't break the pipeline
        
    except Exception as e:
        print(f"[PROCESS ERROR] {filename}: {str(e)}")
        raise


def convert_docx_to_pdf(docx_path: Path, output_pdf: Path) -> bool:
    """
    Convert DOCX to PDF preserving original formatting.
    Uses LibreOffice (cross-platform, Linux-compatible).
    """
    import subprocess
    
    try:
        result = subprocess.run([
            'soffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', str(output_pdf.parent),
            str(docx_path)
        ], check=True, capture_output=True, timeout=60)
        
        # LibreOffice outputs with original name, rename if needed
        libreoffice_output = output_pdf.parent / f"{docx_path.stem}.pdf"
        
        if not libreoffice_output.exists():
            print(f"[PDF ERROR] LibreOffice did not create output file")
            return False
        
        if libreoffice_output != output_pdf:
            libreoffice_output.rename(output_pdf)
        
        print(f"[PDF] Converted via LibreOffice: {docx_path.name} → {output_pdf.name}")
        return True
        
    except subprocess.TimeoutExpired:
        print(f"[PDF ERROR] LibreOffice conversion timed out for {docx_path.name}")
        return False
        
    except (subprocess.CalledProcessError, FileNotFoundError) as e:
        print(f"[PDF ERROR] LibreOffice conversion failed: {e}")
        if hasattr(e, 'stderr') and e.stderr:
            print(f"[PDF ERROR] Stderr: {e.stderr.decode()}")
        return False
        
    except Exception as e:
        print(f"[PDF ERROR] Failed to convert {docx_path}: {e}")
        return False


def create_and_save_metadata_with_pdf(file_path: Path, dst_md: Path, pdf_path: Path, 
                                      total_pages: int, client_id: str, content: str):
    """
    Create metadata with separate PDF view file (preserving original DOCX formatting).
    """
    stat = file_path.stat()
    file_size_mb = round(stat.st_size / 1_000_000, 2)
    updated_date = datetime.fromtimestamp(stat.st_mtime).strftime("%b %Y")
    db = SessionLocal()

    doc_type = detect_doc_type(file_path, content)
    title = create_clean_title(file_path.stem)
    
    # Upload original DOCX (for downloading)
    download_url = upload_original_document(file_path)
    
    # Upload converted PDF (for viewing - preserves original formatting)
    if pdf_path.exists():
        pdf_s3_key = f"{S3_PREFIX_ORIGINAL}{pdf_path.name}"
        view_url = upload_to_s3(pdf_path, pdf_s3_key)
    else:
        view_url = download_url
    
    # Upload parsed markdown (for chunking/search)
    parsed_s3_url = upload_parsed_markdown(dst_md)
    
    # Generate document ID
    existing = db.query(DocumentMetadata).filter(
        DocumentMetadata.client_id == client_id,
        DocumentMetadata.markdown_filename == dst_md.name
    ).first()

    document_id = existing.document_id if existing else str(uuid.uuid4())
    
    metadata_entry = {
        "document_id": document_id,
        "title": title,
        "original_filename": file_path.name,
        "type": doc_type,
        "original_s3_url": view_url,      # PDF preserving original DOCX formatting
        "download_s3_url": download_url,  # Original DOCX
        "parsed_s3_url": parsed_s3_url,   # Markdown for chunking
        "date_updated": updated_date,
        "total_pages": total_pages,
        "file_size": file_size_mb
    }
    
    save_document_metadata(dst_md.name, metadata_entry, client_id)
    db.close()
    
    print(f"[METADATA] Saved for {dst_md.name}")
    print(f"  View: {view_url} (PDF with original formatting)")
    print(f"  Download: {download_url} (Original DOCX)")


def convert_pptx_to_pdf(pptx_path: Path, output_pdf: Path) -> bool:
    """
    Convert PPTX/PPT to PDF preserving original formatting.
    Uses LibreOffice (cross-platform, Linux-compatible).
    """
    import subprocess
    
    try:
        result = subprocess.run([
            'soffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', str(output_pdf.parent),
            str(pptx_path)
        ], check=True, capture_output=True, timeout=60)
        
        # LibreOffice outputs with original name, rename if needed
        libreoffice_output = output_pdf.parent / f"{pptx_path.stem}.pdf"
        
        if not libreoffice_output.exists():
            print(f"[PDF ERROR] LibreOffice did not create output file")
            return False
        
        if libreoffice_output != output_pdf:
            libreoffice_output.rename(output_pdf)
        
        print(f"[PDF] Converted via LibreOffice: {pptx_path.name} → {output_pdf.name}")
        return True
        
    except subprocess.TimeoutExpired:
        print(f"[PDF ERROR] LibreOffice conversion timed out for {pptx_path.name}")
        return False
        
    except (subprocess.CalledProcessError, FileNotFoundError) as e:
        print(f"[PDF ERROR] LibreOffice conversion failed: {e}")
        if hasattr(e, 'stderr') and e.stderr:
            print(f"[PDF ERROR] Stderr: {e.stderr.decode()}")
        return False
        
    except Exception as e:
        print(f"[PDF ERROR] Failed to convert {pptx_path}: {e}")
        return False


def convert_ppt_to_pptx(ppt_path: Path, pptx_path: Path) -> bool:
    """Convert .ppt to .pptx using LibreOffice"""
    import subprocess
    try:
        subprocess.run([
            'soffice',
            '--headless',
            '--convert-to', 'pptx',
            '--outdir', str(pptx_path.parent),
            str(ppt_path)
        ], check=True, capture_output=True, timeout=60)
        return pptx_path.exists()
    except Exception as e:
        print(f"[PPT CONVERT ERROR] {e}")
        return False
    

def convert_doc_to_docx(doc_path: Path, docx_path: Path) -> bool:
    """Convert .doc to .docx using LibreOffice"""
    import subprocess
    try:
        subprocess.run([
            'soffice',
            '--headless',
            '--convert-to', 'docx',
            '--outdir', str(docx_path.parent),
            str(doc_path)
        ], check=True, capture_output=True, timeout=30)
        return docx_path.exists()
    except Exception as e:
        print(f"[DOC CONVERT ERROR] {e}")
        return False
    

def handle_docx(file_path, client_id: str):
    print(f"[WORD] Client {client_id}: Processing {file_path}")

    if file_path.suffix.lower() == '.doc':
        docx_path = file_path.parent / f"{file_path.stem}.docx"
        if not convert_doc_to_docx(file_path, docx_path):
            print(f"[ERROR] Could not convert .doc to .docx")
            return
        file_path = docx_path

    client_parsed_dir = PARSED_DIR / client_id
    client_parsed_dir.mkdir(parents=True, exist_ok=True)
    
    # We still create markdown for parsing/chunking
    dst_md = client_parsed_dir / f"{file_path.stem}.md"
    
    # Parse document
    doc = Document(file_path)
    
    # Create markdown for chunking (as before)
    md_lines = []
    page_num = 1
    para_count = 0

    for paragraph in doc.paragraphs:
        if para_count % 30 == 0:
            md_lines.append(f"\n<!-- PAGE {page_num} -->\n")
            page_num += 1
        para_count += 1
        
        style = paragraph.style.name if paragraph.style else ""
        text = "".join(docx_md_helper(run) for run in paragraph.runs).strip()

        if style.startswith("Heading"):
            level = int(style.split()[-1]) if style.split()[-1].isdigit() else 1
            md_lines.append(f"{'#' * level} {text}")
        elif style in ("List Bullet", "List Paragraph"):
            md_lines.append(f"- {text}")
        elif "List Number" in style:
            md_lines.append(f"1. {text}")
        elif text:
            md_lines.append(text)
        md_lines.append("")

    md_content = "\n".join(md_lines)
    dst_md.write_text(md_content, encoding="utf-8")
    print(f"[SAVED] Markdown: {dst_md}")

    # Convert DOCX to PDF for viewing (preserving original formatting)
    pdf_path = client_parsed_dir / f"{file_path.stem}_view.pdf"
    convert_docx_to_pdf(file_path, pdf_path)
    
    # Create metadata with both files
    create_and_save_metadata_with_pdf(
        file_path, 
        dst_md, 
        pdf_path,
        len(doc.paragraphs), 
        client_id, 
        md_content
    )

def handle_powerpoint(file_path, client_id: str):
    print(f"[POWERPOINT] Client {client_id}: Processing {file_path}")

    # Convert .ppt to .pptx first if needed
    if file_path.suffix.lower() == '.ppt':
        pptx_path = file_path.parent / f"{file_path.stem}.pptx"
        if not convert_ppt_to_pptx(file_path, pptx_path):
            print(f"[ERROR] Could not convert .ppt to .pptx")
            return
        file_path = pptx_path

    client_parsed_dir = PARSED_DIR / client_id
    client_parsed_dir.mkdir(parents=True, exist_ok=True)
    
    # Create markdown for parsing/chunking
    dst_md = client_parsed_dir / f"{file_path.stem}.md"
    
    # Parse presentation
    pres = Presentation(file_path)
    md_lines = []

    for idx, slide in enumerate(pres.slides, start=1):
        md_lines.append(f"\n<!-- PAGE {idx} -->\n")
        md_lines.append(f"# Slide {idx}")
        md_lines.append("")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                md_lines.append(shape.text.strip())
                md_lines.append("")
        md_lines.append("")

    md_content = "\n".join(md_lines)
    dst_md.write_text(md_content, encoding="utf-8")
    print(f"[SAVED] Markdown: {dst_md}")

    # Convert PPTX to PDF for viewing (preserving original formatting)
    pdf_path = client_parsed_dir / f"{file_path.stem}_view.pdf"
    convert_pptx_to_pdf(file_path, pdf_path)
    
    # Create metadata with both files
    create_and_save_metadata_with_pdf(
        file_path, 
        dst_md, 
        pdf_path,
        len(pres.slides), 
        client_id, 
        md_content
    )

def handle_markdown(file_path, client_id: str):
    print(f"[MARKDOWN] Client {client_id}: Processing {file_path}")
    
    # Save to client-specific directory
    client_parsed_dir = PARSED_DIR / client_id
    client_parsed_dir.mkdir(parents=True, exist_ok=True)
    dst = client_parsed_dir / file_path.name
    
    md_content = file_path.read_text(encoding='utf-8')
    lines = md_content.splitlines()
    
    page_lines = []
    page_num = 1
    for i, line in enumerate(lines):
        if i % 50 == 0:
            page_lines.append(f"\n<!-- PAGE {page_num} -->\n")
            page_num += 1
        page_lines.append(line)
    
    md_content = "\n".join(page_lines)
    dst.write_text(md_content, encoding='utf-8')
    
    line_count = len(lines)
    create_and_save_metadata(file_path, dst, line_count, client_id, md_content)

def handle_excel(file_path, client_id: str):
    print(f"[EXCEL] Client {client_id}: Processing {file_path}")
    
    # Save to client-specific directory
    client_parsed_dir = PARSED_DIR / client_id
    client_parsed_dir.mkdir(parents=True, exist_ok=True)
    
    # Convert to CSV (all sheets combined or first sheet only)
    dst = client_parsed_dir / f"{file_path.stem}.csv"
    
    if dst.exists():
        print(f"[SKIP] CSV already exists → {dst}")
        try:
            import csv
            with open(dst, 'r', encoding='utf-8') as f:
                row_count = sum(1 for _ in csv.reader(f)) - 1
            ensure_metadata_exists(file_path, dst, row_count, client_id)
        except Exception as e:
            print(f"[WARNING] Could not check metadata: {e}")
        return
    
    # Convert Excel to CSV
    try:
        excel_file = pd.ExcelFile(file_path)
        
        # If multiple sheets, combine them or use first sheet
        if len(excel_file.sheet_names) == 1:
            df = pd.read_excel(file_path, sheet_name=excel_file.sheet_names[0])
        else:
            # Combine all sheets with sheet name as prefix column
            dfs = []
            for sheet_name in excel_file.sheet_names:
                sheet_df = pd.read_excel(file_path, sheet_name=sheet_name)
                sheet_df.insert(0, 'Sheet', sheet_name)
                dfs.append(sheet_df)
            df = pd.concat(dfs, ignore_index=True)
        
        # Save as CSV
        df.to_csv(dst, index=False, encoding='utf-8')
        print(f"[CONVERTED] {file_path.name} → {dst.name}")
        
        # Count rows
        row_count = len(df)
        
        # Read CSV content for metadata
        csv_content = dst.read_text(encoding='utf-8')
        
        create_and_save_metadata(file_path, dst, row_count, client_id, csv_content)
        
    except Exception as e:
        print(f"[ERROR] Could not convert Excel to CSV: {e}")

def handle_csv(file_path, client_id: str):
    print(f"[CSV] Client {client_id}: Processing {file_path}")
    
    # Save to client-specific directory
    client_parsed_dir = PARSED_DIR / client_id
    client_parsed_dir.mkdir(parents=True, exist_ok=True)
    dst = client_parsed_dir / file_path.name
    
    if dst.exists():
        print(f"[SKIP] File already exists → {dst}")
        try:
            import csv
            with open(file_path, 'r', encoding='utf-8') as f:
                row_count = sum(1 for _ in csv.reader(f)) - 1
            ensure_metadata_exists(file_path, dst, row_count, client_id)
        except Exception as e:
            print(f"[WARNING] Could not check metadata: {e}")
        return
    
    shutil.copy(file_path, dst)
    print(f"[COPIED] {dst}")
    
    # Count rows for metadata
    try:
        import csv
        with open(file_path, 'r', encoding='utf-8') as f:
            row_count = sum(1 for _ in csv.reader(f)) - 1  # -1 for header
        
        create_and_save_metadata(file_path, dst, row_count, client_id)
    except Exception as e:
        print(f"[WARNING] Could not extract CSV metadata: {e}")


def handle_text(file_path, client_id: str):
    print(f"[TEXT] Client {client_id}: Processing {file_path}")
    
    # Save to client-specific directory
    client_parsed_dir = PARSED_DIR / client_id
    client_parsed_dir.mkdir(parents=True, exist_ok=True)
    dst = client_parsed_dir / f"{file_path.stem}.md"
    
    if dst.exists():
        print(f"[SKIP] Markdown already exists → {dst}")
        txt_content = file_path.read_text(encoding='utf-8')
        lines = txt_content.splitlines()
        ensure_metadata_exists(file_path, dst, len(lines), client_id)
        return
    
    txt_content = file_path.read_text(encoding='utf-8')
    lines = txt_content.splitlines()
    
    page_lines = []
    page_num = 1
    for i, line in enumerate(lines):
        if i % 50 == 0:
            page_lines.append(f"\n<!-- PAGE {page_num} -->\n")
            page_num += 1
        page_lines.append(line)
    
    md_content = "\n".join(page_lines)
    dst.write_text(md_content, encoding='utf-8')
    print(f"[SAVED] {dst}")
    
    line_count = len(lines)
    create_and_save_metadata(file_path, dst, line_count, client_id, md_content)


def handle_unknown(file_path, extension, client_id: str = None):
    print(f"[UNKNOWN] Skipping: {file_path} (extension: {extension})")


# FILE HANDLERS
def get_file_handlers():
    return {
        ".pdf": handle_pdf,
        ".docx": handle_docx,
        ".doc": handle_docx,
        ".md": handle_markdown,
        ".markdown": handle_markdown,
        ".txt": handle_text,

        ".xlsx": handle_excel,
        ".xls": handle_excel,
        ".csv": handle_csv,

        ".pptx": handle_powerpoint,
        ".ppt": handle_powerpoint,
    }


# DIRECTORY PROCESSOR
def process_directory(root_path, client_id: str, recursive=True):
    root = Path(root_path)
    if not root.exists():
        raise FileNotFoundError(f"Directory not found: {root_path}")

    handlers = get_file_handlers()
    pattern = "**/*" if recursive else "*"
    file_count = 0

    for item in root.glob(pattern):
        if item.is_file():
            file_count += 1
            extension = item.suffix.lower()
            handler = handlers.get(extension, handle_unknown)

            try:
                # Pass client_id to handler
                if handler == handle_unknown:
                    handler(item, extension)
                else:
                    handler(item, client_id) 
            except Exception as e:
                print(f"Error processing {item}: {e}")
    
    save_metadata_json_backup()
    print(f"\nClient {client_id}: Processed {file_count} files")
